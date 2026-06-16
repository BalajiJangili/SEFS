"""Text extraction utilities for SEFS."""

from __future__ import annotations

import logging
import re
import unicodedata
from pathlib import Path

import fitz

from config import MAX_TEXT_LENGTH


class TextExtractor:
    """Extract normalized text from supported file formats."""

    def __init__(self, max_text_length: int = MAX_TEXT_LENGTH) -> None:
        self.max_text_length = max_text_length
        self.logger = logging.getLogger(self.__class__.__name__)

    def extract(self, filepath: Path) -> str:
        """Extract text from a file and return a normalized truncated string."""
        path = Path(filepath)
        suffix = path.suffix.lower()

        try:
            if suffix in {".txt", ".md", ".py"}:
                text = self._extract_text_file(path)
            elif suffix == ".pdf":
                text = self._extract_pdf(path)
            elif suffix == ".docx":
                text = self._extract_docx(path)
            elif suffix == ".pptx":
                text = self._extract_pptx(path)
            elif suffix == ".csv":
                text = self._extract_csv(path)
            else:
                self.logger.warning("Unsupported file type for extraction: %s", path)
                return ""

            return self._normalize_text(text)
        except Exception as exc:
            self.logger.warning("Failed to extract %s: %s", path, exc)
            return ""

    def _extract_text_file(self, path: Path) -> str:
        """Extract plain text from UTF-8 with latin-1 fallback."""
        try:
            return path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                return path.read_text(encoding="latin-1")
            except UnicodeDecodeError as exc:
                self.logger.warning("Encoding decode failed for %s: %s", path, exc)
                return ""

    def _extract_pdf(self, path: Path) -> str:
        """Extract text from a PDF by concatenating all page text layers."""
        texts: list[str] = []
        try:
            with fitz.open(path) as document:
                for page in document:
                    texts.append(page.get_text())
        except Exception as exc:
            self.logger.warning("PDF extraction failed for %s: %s", path, exc)
            return ""
        return "\n".join(texts)

    def _extract_docx(self, path: Path) -> str:
        """Extract paragraph text from DOCX."""
        try:
            from docx import Document  # type: ignore
        except Exception as exc:
            self.logger.warning("python-docx not available for %s: %s", path, exc)
            return ""

        try:
            document = Document(str(path))
            texts = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
            return "\n".join(texts)
        except Exception as exc:
            self.logger.warning("DOCX extraction failed for %s: %s", path, exc)
            return ""

    def _extract_pptx(self, path: Path) -> str:
        """Extract shape text from PPTX slides."""
        try:
            from pptx import Presentation  # type: ignore
        except Exception as exc:
            self.logger.warning("python-pptx not available for %s: %s", path, exc)
            return ""

        texts: list[str] = []
        try:
            presentation = Presentation(str(path))
            for slide in presentation.slides:
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text:
                        texts.append(str(text))
            return "\n".join(texts)
        except Exception as exc:
            self.logger.warning("PPTX extraction failed for %s: %s", path, exc)
            return ""

    def _extract_csv(self, path: Path) -> str:
        """Extract CSV content into a compact textual table representation."""
        try:
            import pandas as pd  # type: ignore
        except Exception as exc:
            self.logger.warning("pandas not available for %s: %s", path, exc)
            return self._extract_text_file(path)

        try:
            dataframe = pd.read_csv(path, dtype=str, keep_default_na=False)
            if dataframe.empty:
                return ""
            preview_rows = min(len(dataframe), 200)
            preview = dataframe.head(preview_rows)
            return preview.to_csv(index=False)
        except Exception as exc:
            self.logger.warning("CSV parsing failed for %s: %s", path, exc)
            return self._extract_text_file(path)

    def _normalize_text(self, text: str) -> str:
        """Normalize unicode and whitespace, then truncate for embedding."""
        normalized = unicodedata.normalize("NFKC", text)
        normalized = re.sub(r"\s+", " ", normalized).strip()
        if len(normalized) > self.max_text_length:
            return normalized[: self.max_text_length]
        return normalized
